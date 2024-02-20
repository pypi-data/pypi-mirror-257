#!usr/bin/env python
# -*- coding: utf-8 -*-
import codecs
import base64
from collections import defaultdict, Counter
import copy
import contextlib
import logging
import importlib
import datetime
import hashlib
import os
import random
import re
import json
import shlex
import shutil
import subprocess
import sys
import time
import tempfile
import traceback
import urllib
from xmlrpc.client import ServerProxy
import dateparser
import requests
import pyperclip
import toml
from pptx.dml.color import RGBColor

import docx
from docx import Document
from docx.shared import Inches, Pt as DocxPt
from docx.image.exceptions import UnrecognizedImageError
from PIL import Image

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

from chgksuite.common import (
    get_lastdir,
    get_chgksuite_dir,
    set_lastdir,
    DummyLogger,
    log_wrap,
    QUESTION_LABELS,
    check_question,
    replace_escaped,
    retry_wrapper_factory,
    compose_4s,
    tryint,
    custom_csv_to_results,
)
import chgksuite.typotools as typotools
from chgksuite.typotools import (
    remove_excessive_whitespace as rew,
    replace_no_break_spaces,
)

args = None
debug = False
console_mode = False
re_url = re.compile(
    r"""((?:[a-z][\w-]+:(?:/{1,3}|[a-z0-9%])|www\d{0,3}[.]"""
    """|[a-z0-9.\\-]+[.‌​][a-z]{2,4}/)(?:[^\\s<>]+|(([^\\s()<>]+|(([^\\s<>]+)))*))+"""
    """(?:(([^\\s<>]+|(‌​([^\\s<>]+)))*)|[^\\s`!\\[\\]{};:'".,<>?«»“”‘’]))""",
    re.DOTALL,
)
re_perc = re.compile(r"(%[0-9a-fA-F]{2})+")
re_scaps = re.compile(r"(^|[\s])([\[\]\(\)«»А-Я \u0301`ЁA-Z]{2,})([\s,!\.;:-\?]|$)")
re_em = re.compile(r"_(.+?)_")
re_lowercase = re.compile(r"[а-яё]")
re_uppercase = re.compile(r"[А-ЯЁ]")
re_editors = re.compile(r"^[рР]едакторы? *(пакета|тура)? *[—\-–−:] ?")

REQUIRED_LABELS = set(["question", "answer"])
IMGUR_CLIENT_ID = "e86275b3316c6d6"
OVERRIDE_PREFIX = "!!"

ENC = sys.stdout.encoding or "utf8"
CONSOLE_ENC = ENC

WHITEN = {
    "handout": False,
    "zachet": True,
    "nezachet": True,
    "comment": True,
    "source": True,
    "author": False,
}


logger = DummyLogger()
retry_wrapper = None


def unquote(bytestring):
    return urllib.parse.unquote(bytestring.decode("utf8")).encode("utf8")


def make_filename(s, ext, args, addsuffix=""):
    bn = os.path.splitext(os.path.basename(s))[0]
    if addsuffix:
        bn += addsuffix
    if not args.add_ts:
        return bn + "." + ext
    return "{}_{}.{}".format(bn, datetime.datetime.now().strftime("%Y%m%dT%H%M"), ext)


@contextlib.contextmanager
def make_temp_directory(**kwargs):
    temp_dir = tempfile.mkdtemp(**kwargs)
    yield temp_dir
    shutil.rmtree(temp_dir)


def proportional_resize(tup):
    if max(tup) > 600:
        return tuple([int(x * 600 / max(tup)) for x in tup])
    if max(tup) < 200:
        return tuple([int(x * 200 / max(tup)) for x in tup])
    return tup


def imgsize(imgfile):
    img = Image.open(imgfile)
    width, height = proportional_resize((img.width, img.height))
    return width, height


def convert_size(width, height, dimensions="pixels", emsize=25, dpi=120):
    if dimensions == "pixels":
        return width, height
    if dimensions == "ems":
        return width / emsize, height / emsize
    if dimensions == "inches":
        return width / dpi, height / dpi


def search_for_imgfile(imgfile, tmp_dir, targetdir):
    if os.path.isfile(imgfile):
        return imgfile
    for dirname in [tmp_dir, targetdir]:
        if not os.path.isdir(dirname):
            continue
        imgfile2 = os.path.join(dirname, os.path.basename(imgfile))
        if os.path.isfile(imgfile2):
            return imgfile2
    raise Exception("Image file {} not found\n".format(imgfile))


def parse_single_size(ssize, dpi=120, emsize=25):
    if ssize.endswith("in"):
        ssize = ssize[:-2]
        return float(ssize) * dpi
    if ssize.endswith("em"):
        ssize = ssize[:-2]
        return float(ssize) * emsize
    if ssize.endswith("px"):
        ssize = ssize[:-2]
    return float(ssize)


def parseimg(s, dimensions="pixels", tmp_dir=None, targetdir=None):
    width = -1
    height = -1
    sp = shlex.split(s)
    imgfile = sp[-1]
    imgfile = search_for_imgfile(imgfile, tmp_dir, targetdir)
    size = imgsize(imgfile)
    if "big" in sp:
        big = True
        sp = [x for x in sp if x != "big"]
    else:
        big = False

    if len(sp) == 1:
        width, height = convert_size(*size, dimensions=dimensions)
    else:
        for spsp in sp[:-1]:
            spspsp = spsp.split("=")
            if spspsp[0] == "w":
                width = parse_single_size(spspsp[1])
            if spspsp[0] == "h":
                height = parse_single_size(spspsp[1])
        if width != -1 and height == -1:
            height = size[1] * (width / size[0])
        elif width == -1 and height != -1:
            width = size[0] * (height / size[1])
        width, height = convert_size(width, height, dimensions=dimensions)
    return {
        "imgfile": imgfile.replace("\\", "/"),
        "width": width,
        "height": height,
        "big": big,
    }


def partition(alist, indices):
    return [alist[i:j] for i, j in zip([0] + indices, indices + [None])]


def starts_either(s, i, variants):
    for v in variants:
        if s[i : i + len(v)] == v:
            return True
    return False


def parse_4s_elem(s):
    def find_next_unescaped(ss, index):
        j = index + 1
        while j < len(ss):
            if ss[j] == "\\" and j + 2 < len(ss):
                j += 2
            if ss[j] == ss[index]:
                return j
            j += 1
        return -1

    s = s.replace("\\_", "$$$$UNDERSCORE$$$$")
    for gr in re_url.finditer(s):
        gr0 = gr.group(0)
        s = s.replace(gr0, gr0.replace("_", "$$$$UNDERSCORE$$$$"))

    grs = sorted(
        [match.group(0) for match in re_perc.finditer(s)], key=len, reverse=True
    )
    for gr in grs:
        try:
            s = s.replace(gr, unquote(gr.encode("utf8")).decode("utf8"))
        except Exception as e:
            logger.debug(f"error decoding on line {log_wrap(gr)}: {type(e)} {e}\n")

    i = 0
    topart = []
    while i < len(s):
        if s[i] == "_" and (i == 0 or s[i - 1] not in {"\\", "\u6565"}):
            logger.debug("found _ at {} of line {}".format(i, s))
            topart.append(i)
            if find_next_unescaped(s, i) != -1:
                topart.append(find_next_unescaped(s, i) + 1)
                i = find_next_unescaped(s, i) + 2
                continue
        if (
            s[i] == "("
            and i + len("(img") < len(s)
            and "".join(s[i : i + len("(img")]) == "(img"
        ):
            topart.append(i)
            if typotools.find_matching_closing_bracket(s, i) is not None:
                topart.append(typotools.find_matching_closing_bracket(s, i) + 1)
                i = typotools.find_matching_closing_bracket(s, i)
        if (
            s[i] == "("
            and i + len("(screen") < len(s)
            and "".join(s[i : i + len("(screen")]) == "(screen"
        ):
            topart.append(i)
            if typotools.find_matching_closing_bracket(s, i) is not None:
                topart.append(typotools.find_matching_closing_bracket(s, i) + 1)
                i = typotools.find_matching_closing_bracket(s, i)
        if s[i : i + len("(PAGEBREAK)")] == "(PAGEBREAK)":
            topart.append(i)
            topart.append(i + len("(PAGEBREAK)"))
        if starts_either(s, i, ("http://", "https://")):
            topart.append(i)
            j = i + 1
            bracket_level = 0
            while j < len(s) and not (
                s[j].isspace() or s[j] == ")" and bracket_level == 0
            ):
                if s[j] == "(":
                    bracket_level += 1
                elif s[j] == ")" and bracket_level > 0:
                    bracket_level -= 1
                j += 1
            if s[j - 1] in (",", ".", ";"):
                topart.append(j - 1)
            else:
                topart.append(j)
            i = j + 1
        i += 1

    topart = sorted(topart)

    parts = [["", "".join(x.replace("\u6565", ""))] for x in partition(s, topart)]

    def _process(s):
        s = s.replace("\\_", "_")
        s = s.replace("\\.", ".")
        s = s.replace("$$$$UNDERSCORE$$$$", "_")
        return s

    for part in parts:
        if not part[1]:
            continue
        try:
            if part[1][-1] == "_":
                part[1] = part[1][1:]
                part[0] = "em"
            if not part[1]:
                continue
            if part[1][-1] == "_":
                part[1] = part[1][:-1]
                part[0] = "em"
            if not part[1]:
                continue
            if part[1] == "(PAGEBREAK)":
                part[0] = "pagebreak"
                part[1] = ""
            if len(part[1]) > 4 and part[1][:4] == "(img":
                if part[1][-1] != ")":
                    part[1] = part[1] + ")"
                part[1] = part[1][4:-1]
                part[0] = "img"
                logger.debug("found img at {}".format(part[1]))
            if len(part[1]) > 7 and part[1][:7] == "(screen":
                if part[1][-1] != ")":
                    part[1] = part[1] + ")"
                for_print, for_screen = part[1][8:-1].split("|")
                for_print = _process(for_print)
                for_screen = _process(for_screen)
                part[1] = {"for_print": for_print, "for_screen": for_screen}
                part[0] = "screen"
                continue
            if part[1].startswith(("http://", "https://")):
                part[0] = "hyperlink"
            if len(part[1]) > 3 and part[1][:4] == "(sc":
                if part[1][-1] != ")":
                    part[1] = part[1] + ")"
                part[1] = part[1][3:-1]
                part[0] = "sc"
                logger.debug("found img at {}".format(log_wrap(part[1])))
            part[1] = _process(part[1])
        except Exception as e:
            sys.stderr.write(f"Error on part {log_wrap(part)}: {type(e)} {e}")

    return parts


def process_list(element):
    if "-" not in element[1]:
        return
    sp = element[1].split("\n")
    sp = [rew(x) for x in sp]
    list_markers = [i for i in range(len(sp)) if sp[i].startswith("-")]
    if not list_markers:
        return
    preamble = "\n".join(sp[: list_markers[0]])
    inner_list = []
    for num, index in enumerate(list_markers):
        if (num + 1) == len(list_markers):
            inner_list.append(rew("\n".join(sp[index:])[1:]))
        else:
            inner_list.append(rew("\n".join(sp[index : list_markers[num + 1]])[1:]))
    if len(inner_list) == 1:
        element[1] = rew(re.sub("(^|\n)- +", "\\1", element[1]))
    elif preamble:
        element[1] = [preamble, inner_list]
    else:
        element[1] = inner_list


def parse_4s(s, randomize=False):
    mapping = {
        "#": "meta",
        "##": "section",
        "###": "heading",
        "###LJ": "ljheading",
        "#EDITOR": "editor",
        "#DATE": "date",
        "?": "question",
        "№": "number",
        "№№": "setcounter",
        "!": "answer",
        "=": "zachet",
        "!=": "nezachet",
        "^": "source",
        "/": "comment",
        "@": "author",
        ">": "handout",
    }

    structure = []

    if s[0] == "\ufeff" and len(s) > 1:
        s = s[1:]

    with codecs.open("raw.debug", "w", "utf8") as debugf:
        debugf.write(log_wrap(s.split("\n")))

    for line in s.split("\n"):
        if rew(line) == "":
            structure.append(["", ""])
        else:
            if line.split()[0] in mapping:
                structure.append(
                    [mapping[line.split()[0]], rew(line[len(line.split()[0]) :])]
                )
            else:
                if len(structure) >= 1:
                    structure[len(structure) - 1][1] += "\n" + line

    final_structure = []
    current_question = {}
    counter = 1

    if debug:
        with codecs.open("debug1st.debug", "w", "utf8") as debugf:
            debugf.write(log_wrap(structure))

    for element in structure:

        # find list in element

        process_list(element)

        if element[0] in QUESTION_LABELS:
            if element[0] in current_question:

                if isinstance(current_question[element[0]], str) and isinstance(
                    element[1], str
                ):
                    current_question[element[0]] += "\n" + element[1]

                elif isinstance(current_question[element[0]], list) and isinstance(
                    element[1], str
                ):
                    current_question[element[0]][0] += "\n" + element[1]

                elif isinstance(current_question[element[0]], str) and isinstance(
                    element[1], list
                ):
                    current_question[element[0]] = [
                        element[1][0] + "\n" + current_question[element[0]],
                        element[1][1],
                    ]

                elif isinstance(current_question[element[0]], list) and isinstance(
                    element[1], list
                ):
                    current_question[element[0]][0] += "\n" + element[1][0]
                    current_question[element[0]][1] += element[1][1]
            else:
                current_question[element[0]] = element[1]

        elif element[0] == "":

            if current_question != {} and set(current_question.keys()) != {
                "setcounter"
            }:

                try:
                    assert all(
                        (True if label in current_question else False)
                        for label in REQUIRED_LABELS
                    )
                except AssertionError:
                    logger.error(
                        "Question {} misses "
                        "some of the required fields "
                        "and will therefore "
                        "be omitted.".format(log_wrap(current_question))
                    )
                    continue
                if "setcounter" in current_question:
                    counter = int(current_question["setcounter"])
                if "number" not in current_question:
                    current_question["number"] = counter
                    counter += 1
                final_structure.append(["Question", current_question])

                current_question = {}

        else:
            final_structure.append([element[0], element[1]])

    if current_question != {}:
        try:
            assert all(
                (True if label in current_question else False)
                for label in REQUIRED_LABELS
            )
            if "setcounter" in current_question:
                counter = int(current_question["setcounter"])
            if "number" not in current_question:
                current_question["number"] = counter
                counter += 1
            final_structure.append(["Question", current_question])
        except AssertionError:
            logger.error(
                "Question {} misses "
                "some of the required fields and will therefore "
                "be omitted.".format(log_wrap(current_question))
            )

    if randomize:
        random.shuffle(final_structure, lambda: 0.3)
        i = 1
        for element in final_structure:
            if element[0] == "Question":
                element[1]["number"] = i
                i += 1

    if debug:
        with codecs.open("debug.debug", "w", "utf8") as debugf:
            debugf.write(log_wrap(final_structure))

    for element in final_structure:
        if element[0] == "Question":
            check_question(element[1], logger=logger)
            for field in [
                "handout",
                "question",
                "answer",
                "zachet",
                "nezachet",
                "comment",
                "source",
                "author",
            ]:
                val = element[1].get(field)
                if val is None:
                    continue
                is_list = False
                if isinstance(val, list):
                    is_list = True
                    val = val[0]
                sp = val.split(" ", 1)
                if len(sp) == 1:
                    continue
                sp1, sp2 = sp
                if sp1.startswith(OVERRIDE_PREFIX):
                    if "overrides" not in element[1]:
                        element[1]["overrides"] = {}
                    element[1]["overrides"][field] = sp1[
                        len(OVERRIDE_PREFIX) :
                    ].replace("~", " ")
                    if is_list:
                        element[1][field][0] = sp2
                    else:
                        element[1][field] = sp2

    return final_structure


def md5(s):
    return hashlib.md5(s).hexdigest()


def find_heading(structure):
    h_id = -1
    for e, x in enumerate(structure):
        if x[0] == "ljheading":
            return (e, x)
        elif x[0] == "heading":
            h_id = e
    if h_id >= 0:
        return (h_id, structure[h_id])
    return None


def find_tour(structure):
    for e, x in enumerate(structure):
        if x[0] == "section":
            return (e, x)
    return None


def check_if_zero(Question):
    number = Question.get("number")
    if number is None:
        return False
    if isinstance(number, int) and number == 0:
        return True
    if isinstance(number, str) and number.startswith(("0", "Размин")):
        return True
    return False


def gui_compose(largs, sourcedir=None):
    global args
    global console_mode
    args = largs
    global debug
    global logger
    global retry_wrapper
    assert sourcedir is not None

    logger = logging.getLogger("composer")
    logger.setLevel(logging.DEBUG)
    fh = logging.FileHandler("composer.log", encoding="utf8")
    fh.setLevel(logging.DEBUG)
    ch = logging.StreamHandler()
    if args.debug:
        ch.setLevel(logging.DEBUG)
    else:
        ch.setLevel(logging.INFO)
    formatter = logging.Formatter("%(asctime)s | %(message)s")
    fh.setFormatter(formatter)
    ch.setFormatter(formatter)
    logger.addHandler(fh)
    logger.addHandler(ch)

    retry_wrapper = retry_wrapper_factory(logger)

    if args.debug:
        debug = True

    argsdict = vars(args)
    logger.debug(log_wrap(argsdict))

    if args.filename and args.filetype:
        if args.filetype == "lj":
            if args.login and args.password:
                console_mode = True
        else:
            console_mode = True

    ld = get_lastdir()
    if args.filename:
        if isinstance(args.filename, list):
            ld = os.path.dirname(os.path.abspath(args.filename[0]))
        else:
            ld = os.path.dirname(os.path.abspath(args.filename))
    set_lastdir(ld)
    if not args.filename:
        print("No file specified.")
        sys.exit(1)

    if isinstance(args.filename, list):
        if not args.merge:
            for fn in args.filename:
                targetdir = os.path.dirname(os.path.abspath(fn))
                filename = os.path.basename(os.path.abspath(fn))
                process_file_wrapper(filename, sourcedir, targetdir)
        else:
            targetdir = os.path.dirname(os.path.abspath(args.filename[0]))
            process_file_wrapper(args.filename, sourcedir, targetdir)
    else:
        targetdir = os.path.dirname(os.path.abspath(args.filename))
        filename = os.path.basename(os.path.abspath(args.filename))
        process_file_wrapper(filename, sourcedir, targetdir)


def process_file_wrapper(filename, sourcedir, targetdir):
    resourcedir = os.path.join(sourcedir, "resources")
    with make_temp_directory(dir=get_chgksuite_dir()) as tmp_dir:
        for fn in [
            args.docx_template,
            os.path.join(resourcedir, "fix-unnumbered-sections.sty"),
            args.tex_header,
        ]:
            shutil.copy(fn, tmp_dir)
        process_file(filename, tmp_dir, sourcedir, targetdir)


def parse_filepath(filepath):
    with codecs.open(filepath, "r", "utf8") as input_file:
        input_text = input_file.read()
    input_text = input_text.replace("\r", "")
    return parse_4s(input_text, randomize=args.randomize)


def make_merged_filename(filelist):
    filelist = [os.path.splitext(os.path.basename(x))[0] for x in filelist]
    prefix = os.path.commonprefix(filelist)
    suffix = "_".join(x[len(prefix) :] for x in filelist)
    return prefix + suffix


def generate_navigation(strus):
    titles = [x[0][0]["header"].split(". ")[-1] for x in strus]
    urls = [x[1]["url"] for x in strus]
    result = []
    for i in range(len(titles)):
        inner = []
        for j in range(len(urls)):
            inner.append(
                titles[j]
                if j == i
                else '<a href="{}">{}</a>'.format(urls[j], titles[j])
            )
        result.append(" | ".join(inner))
    return result


def find_min_content_index(structure):
    types_ = [x[0] for x in structure]
    try:
        min_section = types_.index("section")
    except ValueError:
        min_section = None
    try:
        min_question = types_.index("Question")
    except ValueError:
        min_question = None
    if min_section is not None and min_question is not None:
        return min(min_section, min_question)
    elif min_section is not None:
        return min_section
    else:
        return min_question


def backtick_replace(el):
    while "`" in el:
        if el.index("`") + 1 >= len(el):
            el = el.replace("`", "")
        else:
            if el.index("`") + 2 < len(el) and re.search(r"\s", el[el.index("`") + 2]):
                el = el[: el.index("`") + 2] + "" + el[el.index("`") + 2 :]
            if el.index("`") + 1 < len(el) and re_lowercase.search(
                el[el.index("`") + 1]
            ):
                el = (
                    el[: el.index("`") + 1]
                    + ""
                    + el[el.index("`") + 1]
                    + "\u0301"
                    + el[el.index("`") + 2 :]
                )
            elif el.index("`") + 1 < len(el) and re_uppercase.search(
                el[el.index("`") + 1]
            ):
                el = (
                    el[: el.index("`") + 1]
                    + ""
                    + el[el.index("`") + 1]
                    + "\u0301"
                    + el[el.index("`") + 2 :]
                )
            el = el[: el.index("`")] + el[el.index("`") + 1 :]
    return el


class BaseExporter:
    def __init__(self, *args, **kwargs):
        self.structure = args[0]
        self.args = args[1]
        self.dir_kwargs = args[2]
        with open(self.args.labels_file, encoding="utf8") as f:
            self.labels = toml.load(f)

    def get_label(self, question, field, number=None):
        if field in ("question", "tour"):
            lbl = (question.get("overrides") or {}).get(field) or self.labels[
                "question_labels"
            ][field]
            num = question.get("number") or number
            if self.args.language in ("uz", "uz_cyr"):
                return f"{num} – {lbl}"
            else:
                return f"{lbl} {num}"
        if field in (question.get("overrides") or {}):
            return question["overrides"][field]
        if field == "source" and isinstance(question.get("source" or ""), list):
            return self.labels["question_labels"]["sources"]
        return self.labels["question_labels"][field]

    def remove_square_brackets(self, s):
        hs = self.labels["question_labels"]["handout_short"]
        s = s.replace("\\[", "LEFTSQUAREBRACKET")
        s = s.replace("\\]", "RIGHTSQUAREBRACKET")
        s = re.sub(f"\\[{hs}(.+?)\\]", "{" + hs + "\\1}", s, flags=re.DOTALL)
        i = 0
        while "[" in s and "]" in s and i < 10:
            s = re.sub(" *\\[.+?\\]", "", s, flags=re.DOTALL)
            s = s.strip()
            i += 1
        if i == 10:
            sys.stderr.write(
                f"Error replacing square brackets on question: {s}, retries exceeded\n"
            )
        s = re.sub("\\{" + hs + "(.+?)\\}", "[" + hs + "\\1]", s, flags=re.DOTALL)
        s = s.replace("LEFTSQUAREBRACKET", "[")
        s = s.replace("RIGHTSQUAREBRACKET", "]")
        return s


class DbExporter(BaseExporter):
    BASE_MAPPING = {
        "section": "Тур",
        "heading": "Чемпионат",
        "editor": "Редактор",
        "meta": "Инфо",
    }
    re_date_sep = re.compile(" [—–-] ")

    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.qcount = 0
        self.im = Imgur(self.args.imgur_client_id or IMGUR_CLIENT_ID)

    def baseyapper(self, e):
        if isinstance(e, str):
            return self.base_element_layout(e)
        elif isinstance(e, list):
            if not any(isinstance(x, list) for x in e):
                return self.base_element_layout(e)
            else:
                return "\n".join([self.base_element_layout(x) for x in e])

    def parse_and_upload_image(self, path):
        parsed_image = parseimg(
            path,
            dimensions="pixels",
            targetdir=self.dir_kwargs.get("targetdir"),
            tmp_dir=self.dir_kwargs.get("tmp_dir"),
        )
        imgfile = parsed_image["imgfile"]
        w = parsed_image["width"]
        h = parsed_image["height"]
        pil_image = Image.open(imgfile)
        w_orig, h_orig = pil_image.size
        if w_orig != w or h_orig != h:
            logger.info("resizing image {}".format(imgfile))
            pil_image = pil_image.resize((int(w), int(h)), resample=Image.LANCZOS)
            bn, ext = os.path.splitext(imgfile)
            resized_fn = "{}_resized.png".format(bn)
            pil_image.save(resized_fn, "PNG")
            to_upload = resized_fn
        else:
            to_upload = imgfile
        logger.info("uploading {}...".format(to_upload))
        uploaded_image = self.im.upload_image(to_upload, title=to_upload)
        imglink = uploaded_image["data"]["link"]
        logger.info("the link for {} is {}...".format(to_upload, imglink))
        return imglink

    def baseformat(self, s):
        res = ""
        for run in parse_4s_elem(s):
            if run[0] in ("", "hyperlink"):
                res += run[1].replace("\n", "\n   ")
            if run[0] == "em":
                res += run[1]
            if run[0] == "screen":
                res += run[1]["for_print"]
            if run[0] == "img":
                if run[1].startswith(("http://", "https://")):
                    imglink = run[1]
                else:
                    imglink = self.parse_and_upload_image(run[1])
                res += "(pic: {})".format(imglink)
        while res.endswith("\n"):
            res = res[:-1]
        res = replace_escaped(res)
        return res

    def base_element_layout(self, e):
        res = ""
        if isinstance(e, str):
            if self.args.remove_accents:
                for match in re.finditer("(.)\u0301", e):
                    replacement = match.group(1).upper()
                    e = e.replace(match.group(0), replacement)
            res = self.baseformat(e)
            return res
        if isinstance(e, list):
            res = "\n".join(
                [
                    "   {}. {}".format(i + 1, self.base_element_layout(x))
                    for i, x in enumerate(e)
                ]
            )
        return res

    @staticmethod
    def wrap_date(s):
        s = s.strip()
        parsed = dateparser.parse(s)
        if isinstance(parsed, datetime.datetime):
            parsed = parsed.date()
        if not parsed:
            logger.error(
                "unable to parse date {}, setting to default 2010-01-01".format(s)
            )
            return datetime.date(2010, 1, 1).strftime("%d-%b-%Y")
        if parsed > datetime.date.today():
            parsed = parsed.replace(year=parsed.year - 1)
        formatted = parsed.strftime("%d-%b-%Y")
        return formatted

    def base_format_element(self, pair):
        if pair[0] == "Question":
            return self.base_format_question(pair[1])
        if pair[0] in self.BASE_MAPPING:
            return "{}:\n{}\n\n".format(
                self.BASE_MAPPING[pair[0]], self.baseyapper(pair[1])
            )
        elif pair[0] == "date":
            re_search = self.re_date_sep.search(pair[1])
            if re_search:
                gr0 = re_search.group(0)
                dates = pair[1].split(gr0)
                return "Дата:\n{} - {}\n\n".format(
                    self.wrap_date(dates[0]), self.wrap_date(dates[-1])
                )
            else:
                return "Дата:\n{}\n\n".format(self.wrap_date(pair[1]))

    @staticmethod
    def _get_last_value(dct, key):
        if isinstance(dct[key], list):
            return dct[key][-1]
        return dct[key]

    @staticmethod
    def _add_to_dct(dct, key, to_add):
        if isinstance(dct[key], list):
            dct[key][-1] += to_add
        else:
            dct[key] += to_add

    def base_format_question(self, q):

        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        res = "Вопрос {}:\n{}\n\n".format(
            self.qcount if "number" not in q else q["number"],
            self.baseyapper(q["question"]),
        )
        if "number" not in q:
            self.qcount += 1
        res += "Ответ:\n{}\n\n".format(self.baseyapper(q["answer"]))
        if "zachet" in q:
            res += "Зачет:\n{}\n\n".format(self.baseyapper(q["zachet"]))
        if "nezachet" in q:
            res += "Незачет:\n{}\n\n".format(self.baseyapper(q["zachet"]))
        if "comment" in q:
            res += "Комментарий:\n{}\n\n".format(self.baseyapper(q["comment"]))
        if "source" in q:
            res += "Источник:\n{}\n\n".format(self.baseyapper(q["source"]))
        if "author" in q:
            res += "Автор:\n{}\n\n".format(self.baseyapper(q["author"]))
        return res

    def export(self, outfilename):

        result = []
        lasttour = 0
        zeroq = 1
        for i, pair in enumerate(self.structure):
            if pair[0] == "section":
                lasttour = i
            while (
                pair[0] == "meta"
                and (i + 1) < len(self.structure)
                and self.structure[i + 1][0] == "meta"
            ):
                pair[1] += "\n{}".format(self.structure[i + 1][1])
                self.structure.pop(i + 1)
            if pair[0] == "Question" and check_if_zero(pair[1]):
                tourheader = "Нулевой вопрос {}".format(zeroq)
                zeroq += 1
                pair[1]["number"] = 1
                self.structure.insert(lasttour, self.structure.pop(i))
                self.structure.insert(lasttour, ["section", tourheader])
        for pair in self.structure:
            if pair[0] == "Question" and "nezachet" in pair[1]:
                field = "zachet" if "zachet" in pair[1] else "answer"
                last_val = self._get_last_value(pair[1], field)
                nezachet = self.baseyapper(pair[1].pop("nezachet"))
                to_add = "{}\n   Незачёт: {}".format(
                    "." if not last_val.endswith(".") else "", nezachet
                )
                self._add_to_dct(pair[1], field, to_add)
            if pair[0] == "editor":
                pair[1] = re.sub(re_editors, "", pair[1])
                logger.info('Поле "Редактор" было автоматически изменено.')
            res = self.base_format_element(pair)
            if res:
                result.append(res)
        text = "".join(result)
        with codecs.open(outfilename, "w", "utf8") as f:
            f.write(text)
        logger.info("Output: {}".format(outfilename))
        if self.args.clipboard:
            pyperclip.copy(text)


class RedditExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.im = Imgur(self.args.imgur_client_id or IMGUR_CLIENT_ID)
        self.qcount = 1

    def reddityapper(self, e):
        if isinstance(e, str):
            return self.reddit_element_layout(e)
        elif isinstance(e, list):
            if not any(isinstance(x, list) for x in e):
                return self.reddit_element_layout(e)
            else:
                return "  \n".join([self.reddit_element_layout(x) for x in e])

    def parse_and_upload_image(self, path):
        parsed_image = parseimg(
            path,
            dimensions="ems",
            targetdir=self.dir_kwargs.get("targetdir"),
            tmp_dir=self.dir_kwargs.get("tmp_dir"),
        )
        imgfile = parsed_image["imgfile"]
        if os.path.isfile(imgfile):
            uploaded_image = self.im.upload_image(imgfile, title=imgfile)
            imglink = uploaded_image["data"]["link"]
            return imglink

    def redditformat(self, s):
        res = ""
        for run in parse_4s_elem(s):
            if run[0] in ("", "hyperlink"):
                res += run[1]
            if run[0] == "screen":
                res += run[1]["for_screen"]
            if run[0] == "em":
                res += "_{}_".format(run[1])
            if run[0] == "img":
                if run[1].startswith(("http://", "https://")):
                    imglink = run[1]
                else:
                    imglink = self.parse_and_upload_image(run[1])
                res += "[картинка]({})".format(imglink)
        while res.endswith("\n"):
            res = res[:-1]
        res = res.replace("\n", "  \n")
        return res

    def reddit_element_layout(self, e):
        res = ""
        if isinstance(e, str):
            res = self.redditformat(e)
            return res
        if isinstance(e, list):
            res = "  \n".join(
                [
                    "{}\\. {}".format(i + 1, self.reddit_element_layout(x))
                    for i, x in enumerate(e)
                ]
            )
        return res

    def reddit_format_element(self, pair):
        if pair[0] == "Question":
            return self.reddit_format_question(pair[1])

    def reddit_format_question(self, q):
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        res = "__Вопрос {}__: {}  \n".format(
            self.qcount if "number" not in q else q["number"],
            self.reddityapper(q["question"]),
        )
        if "number" not in q:
            self.qcount += 1
        res += "__Ответ:__ >!{}  \n".format(self.reddityapper(q["answer"]))
        if "zachet" in q:
            res += "__Зачёт:__ {}  \n".format(self.reddityapper(q["zachet"]))
        if "nezachet" in q:
            res += "__Незачёт:__ {}  \n".format(self.reddityapper(q["nezachet"]))
        if "comment" in q:
            res += "__Комментарий:__ {}  \n".format(self.reddityapper(q["comment"]))
        if "source" in q:
            res += "__Источник:__ {}  \n".format(self.reddityapper(q["source"]))
        if "author" in q:
            res += "!<\n__Автор:__ {}  \n".format(self.reddityapper(q["author"]))
        else:
            res += "!<\n"
        return res

    def export(self, outfile):
        result = []
        for pair in self.structure:
            res = self.reddit_format_element(pair)
            if res:
                result.append(res)
        text = "\n\n".join(result)
        with codecs.open(outfile, "w", "utf8") as f:
            f.write(text)
        logger.info("Output: {}".format(outfile))


class TelegramExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.chgksuite_dir = get_chgksuite_dir()
        self.pyrogram = importlib.import_module(
            "pyrogram"
        )  # pyrogram slows down startup quite a bit, so only import it when needed
        try:
            self.init_tg()
        except self.pyrogram.errors.exceptions.unauthorized_401.AuthKeyUnregistered:
            filepath = os.path.join(
                self.chgksuite_dir, self.args.tgaccount + ".session"
            )
            if os.path.isfile(filepath):
                os.remove(filepath)
            self.init_tg()
        self.qcount = 1
        self.number = 1
        self.tg_heading = None

    def init_tg(self):
        api_id, api_hash = self.get_api_credentials()
        self.app = self.pyrogram.Client(
            self.args.tgaccount,
            api_id,
            api_hash,
            workdir=self.chgksuite_dir,
            hide_password=not self.args.no_hide_password,
        )
        with self.app:
            logger.debug(self.app.get_me())

    def get_api_credentials(self):
        pyrogram_toml_file_path = os.path.join(self.chgksuite_dir, "pyrogram.toml")
        if os.path.exists(pyrogram_toml_file_path) and not self.args.reset_api:
            with open(pyrogram_toml_file_path, "r", encoding="utf8") as f:
                pyr = toml.load(f)
            return pyr["api_id"], pyr["api_hash"]
        else:
            print("Please enter you api_id and api_hash.")
            print(
                "Go to https://my.telegram.org/apps, register an app and paste the credentials here."
            )
            api_id = input("Enter your api_id: ").strip()
            api_hash = input("Enter your api_hash: ").strip()
            with open(pyrogram_toml_file_path, "w", encoding="utf8") as f:
                toml.dump({"api_id": api_id, "api_hash": api_hash}, f)
            return api_id, api_hash

    def tgyapper(self, e):
        if isinstance(e, str):
            return self.tg_element_layout(e)
        elif isinstance(e, list):
            if not any(isinstance(x, list) for x in e):
                return self.tg_element_layout(e)
            else:
                res = []
                images = []
                for x in e:
                    res_, images_ = self.tg_element_layout(x)
                    images.extend(images_)
                    res.append(res_)
                return "\n".join(res), images

    def tgformat(self, s):
        res = ""
        image = None
        for run in parse_4s_elem(s):
            if run[0] in ("", "hyperlink"):
                res += run[1]
            if run[0] == "screen":
                res += run[1]["for_screen"]
            if run[0] == "em":
                res += "_{}_".format(run[1])
            if run[0] == "img":
                if run[1].startswith(("http://", "https://")):
                    res += run[1]
                else:
                    res += self.labels["general"].get("cf_image", "см. изображение")
                    parsed_image = parseimg(
                        run[1],
                        dimensions="ems",
                        targetdir=self.dir_kwargs.get("targetdir"),
                        tmp_dir=self.dir_kwargs.get("tmp_dir"),
                    )
                    imgfile = parsed_image["imgfile"]
                    if os.path.isfile(imgfile):
                        image = imgfile
                    else:
                        raise Exception(f"image {run[1]} doesn't exist")
        while res.endswith("\n"):
            res = res[:-1]
        if "*" in res and not self.args.disable_asterisks_processing:
            res = res.replace("*", "&#42;")
        return res, image

    def tg_element_layout(self, e):
        res = ""
        images = []
        if isinstance(e, str):
            res, image = self.tgformat(e)
            if image:
                images.append(image)
            return res, images
        if isinstance(e, list):
            result = []
            for i, x in enumerate(e):
                res_, images_ = self.tg_element_layout(x)
                images.extend(images_)
                result.append("{}. {}".format(i + 1, res_))
            res = "\n".join(result)
        return res, images

    def _post(self, chat_id, text, photo, reply_to_message_id=None):
        if photo:
            if not text:
                caption = ""
            elif text == "---":
                caption = "--"
            else:
                caption = "---"
            msg = self.app.send_photo(
                chat_id,
                photo,
                caption=caption,
                parse_mode=self.pyrogram.enums.ParseMode.MARKDOWN,
                reply_to_message_id=reply_to_message_id,
            )
            if text:
                time.sleep(2)
                self.app.edit_message_text(
                    chat_id,
                    msg.id,
                    text=text,
                    parse_mode=self.pyrogram.enums.ParseMode.MARKDOWN,
                    disable_web_page_preview=True,
                )
        else:
            msg = self.app.send_message(
                chat_id,
                text,
                parse_mode=self.pyrogram.enums.ParseMode.MARKDOWN,
                disable_web_page_preview=True,
                reply_to_message_id=reply_to_message_id,
            )
        return msg

    def __post(self, *args, **kwargs):
        retries = 0
        while retries <= 2:
            try:
                return self._post(*args, **kwargs)
            except self.pyrogram.errors.exceptions.flood_420.FloodWait as e:
                mstr = str(e)
                secs_to_wait = re.search("([0-9]+) seconds is required", mstr)
                if secs_to_wait:
                    secs_to_wait = int(secs_to_wait.group(1)) + 30
                else:
                    secs_to_wait = 120
                logger.error(
                    f"Telegram thinks we are spammers, waiting for {secs_to_wait} seconds"
                )
                time.sleep(secs_to_wait)
                retries += 1

    def post(self, posts):
        if self.args.dry_run:
            logger.info("skipping posting due to dry run")
            return
        messages = []
        text, im = posts[0]
        root_msg = self.__post(
            self.channel_id,
            self.labels["general"]["handout_for_question"].format(text[3:])
            if text.startswith("QQQ")
            else text,
            im,
        )
        if (
            len(posts) >= 2 and text.startswith("QQQ") and im and posts[1][0]
        ):  # crutch for case when the question doesn't fit without image
            prev_root_msg = root_msg
            root_msg = self.__post(self.channel_id, posts[1][0], posts[1][1])
            posts = posts[1:]
            messages.append(root_msg)
            messages.append(prev_root_msg)
        time.sleep(2.1)
        root_msg_in_chat = self.app.get_discussion_message(self.channel_id, root_msg.id)
        logger.info(f"Posted message {root_msg.link} ({root_msg_in_chat.link} in chat)")
        time.sleep(random.randint(5, 7))
        if root_msg not in messages:
            messages.append(root_msg)
        messages.append(root_msg_in_chat)
        for post in posts[1:]:
            text, im = post
            reply_msg = self.__post(
                self.chat_id, text, im, reply_to_message_id=root_msg_in_chat.id
            )
            logger.info(
                f"Replied to message {root_msg_in_chat.link} with {reply_msg.link}"
            )
            time.sleep(random.randint(5, 7))
            messages.append(reply_msg)
        return messages

    def post_wrapper(self, posts):
        messages = self.post(posts)
        if self.section and not self.args.dry_run:
            self.section_links.append(messages[0].link)
        self.section = False

    def tg_process_element(self, pair):
        if pair[0] == "Question":
            q = pair[1]
            if "setcounter" in q:
                self.qcount = int(q["setcounter"])
            number = self.qcount if "number" not in q else q["number"]
            self.qcount += 1
            self.number = number
            if self.args.skip_until and (
                not tryint(number) or tryint(number) < self.args.skip_until
            ):
                logger.info(f"skipping question {number}")
                return
            if self.buffer_texts or self.buffer_images:
                posts = self.split_to_messages(self.buffer_texts, self.buffer_images)
                self.post_wrapper(posts)
                self.buffer_texts = []
                self.buffer_images = []
            posts = self.tg_format_question(pair[1], number=number)
            self.post_wrapper(posts)
        elif self.args.skip_until and self.number < self.args.skip_until:
            logger.info(f"skipping element {pair[0]}")
            return
        elif pair[0] == "heading":
            text, images = self.tg_element_layout(pair[1])
            if not self.tg_heading:
                self.tg_heading = text
            self.buffer_texts.append(f"**{text}**")
            self.buffer_images.extend(images)
        elif pair[0] == "section":
            if self.buffer_texts or self.buffer_images:
                posts = self.split_to_messages(self.buffer_texts, self.buffer_images)
                self.post_wrapper(posts)
                self.buffer_texts = []
                self.buffer_images = []
            text, images = self.tg_element_layout(pair[1])
            self.buffer_texts.append(f"**{text}**")
            self.buffer_images.extend(images)
            self.section = True
        else:
            text, images = self.tg_element_layout(pair[1])
            if text:
                self.buffer_texts.append(text)
            if images:
                self.buffer_images.extend(images)

    def assemble(self, list_, lb_after_first=False):
        list_ = [x for x in list_ if x]
        list_ = [x.strip() for x in list_ if not x.startswith("\n||")]
        if lb_after_first:
            list_[0] = list_[0] + "\n"
        res = "\n".join(list_)
        res = res.replace("\n||\n", "\n||")
        while res.endswith("\n"):
            res = res[:-1]
        if res.endswith("\n||"):
            res = res[:-3] + "||"
        if self.args.nospoilers:
            res = res.replace("||", "")
        return res

    def make_chunk(self, texts, images):
        if images:
            im, images = images[0], images[1:]
            threshold = 1024
        else:
            im = None
            threshold = 2048
        if not texts:
            return "", im, texts, images
        if len(texts[0]) <= threshold:
            for i in range(0, len(texts)):
                if i:
                    text = self.assemble(texts[:-i])
                else:
                    text = self.assemble(texts)
                if len(text) <= threshold:
                    if i:
                        texts = texts[-i:]
                    else:
                        texts = []
                    return text, im, texts, images
        else:
            threshold_ = threshold - 3
            chunk = texts[0][:threshold_]
            rest = texts[0][threshold_:]
            if texts[0].endswith("||"):
                chunk += "||"
                rest = "||" + rest
            texts[0] = rest
            return chunk, im, texts, images

    def split_to_messages(self, texts, images):
        result = []
        while texts or images:
            chunk, im, texts, images = self.make_chunk(texts, images)
            if chunk or im:
                result.append((chunk, im))
        return result

    def swrap(self, s_, t="both"):
        if not s_:
            res = s_
        if self.args.nospoilers:
            res = s_
        elif t == "both":
            res = "||" + s_ + "||"
        elif t == "left":
            res = "||" + s_
        elif t == "right":
            res = s_ + "||"
        return res

    @staticmethod
    def lwrap(l_, lb_after_first=False):
        l_ = [x.strip() for x in l_ if x]
        if lb_after_first:
            return l_[0] + "\n" + "\n".join([x for x in l_[1:]])
        return "\n".join(l_)

    def tg_format_question(self, q, number=None):
        txt_q, images_q = self.tgyapper(q["question"])
        txt_q = "**{}:** {}  \n".format(
            self.get_label(q, "question", number=number),
            txt_q,
        )
        if "number" not in q:
            self.qcount += 1
        images_a = []
        txt_a, images_ = self.tgyapper(q["answer"])
        images_a.extend(images_)
        txt_a = "**{}:** {}".format(self.get_label(q, "answer"), txt_a)
        txt_z = ""
        txt_nz = ""
        txt_comm = ""
        txt_s = ""
        txt_au = ""
        if "zachet" in q:
            txt_z, images_ = self.tgyapper(q["zachet"])
            images_a.extend(images_)
            txt_z = "**{}:** {}".format(self.get_label(q, "zachet"), txt_z)
        if "nezachet" in q:
            txt_nz, images_ = self.tgyapper(q["nezachet"])
            images_a.extend(images_)
            txt_nz = "**{}:** {}".format(self.get_label(q, "nezachet"), txt_nz)
        if "comment" in q:
            txt_comm, images_ = self.tgyapper(q["comment"])
            images_a.extend(images_)
            txt_comm = "**{}:** {}".format(self.get_label(q, "comment"), txt_comm)
        if "source" in q:
            txt_s, images_ = self.tgyapper(q["source"])
            images_a.extend(images_)
            txt_s = f"**{self.get_label(q, 'source')}:** {txt_s}"
        if "author" in q:
            txt_au, images_ = self.tgyapper(q["author"])
            images_a.extend(images_)
            txt_au = f"**{self.get_label(q, 'author')}:** {txt_au}"
        q_threshold = 2048 if not images_q else 1024
        full_question = self.assemble(
            [
                txt_q,
                self.swrap(txt_a, t="left"),
                txt_z,
                txt_nz,
                txt_comm,
                self.swrap(txt_s, t="right"),
                txt_au,
            ],
            lb_after_first=True,
        )
        if len(full_question) <= q_threshold:
            res = [(full_question, images_q[0] if images_q else None)]
            for i in images_a:
                res.append(("", i))
            return res
        elif images_q and len(full_question) <= 2048:
            full_question = re.sub(
                "\\[" + self.labels["question_labels"]["handout"] + ": +?\\]\n",
                "",
                full_question,
            )
            res = [(f"QQQ{number}", images_q[0]), (full_question, None)]
            for i in images_a:
                res.append(("", i))
            return res
        q_without_s = self.assemble(
            [
                txt_q,
                self.swrap(txt_a, t="left"),
                txt_z,
                txt_nz,
                self.swrap(txt_comm, t="right"),
            ],
            lb_after_first=True,
        )
        if len(q_without_s) <= q_threshold:
            res = [(q_without_s, images_q[0] if images_q else None)]
            res.extend(
                self.split_to_messages(
                    self.lwrap([self.swrap(txt_s), txt_au]), images_a
                )
            )
            return res
        q_a_only = self.assemble([txt_q, self.swrap(txt_a)], lb_after_first=True)
        if len(q_a_only) <= q_threshold:
            res = [(q_a_only, images_q[0] if images_q else None)]
            res.extend(
                self.split_to_messages(
                    self.lwrap(
                        [
                            self.swrap(txt_z),
                            self.swrap(txt_nz),
                            self.swrap(txt_comm),
                            self.swrap(txt_s),
                            txt_au,
                        ]
                    ),
                    images_a,
                )
            )
            return res
        return self.split_to_messages(
            self.lwrap(
                [
                    txt_q,
                    self.swrap(txt_a),
                    self.swrap(txt_z),
                    self.swrap(txt_nz),
                    self.swrap(txt_comm),
                    self.swrap(txt_s),
                    txt_au,
                ],
                lb_after_first=True,
            ),
            (images_q or []) + (images_a or []),
        )

    def export(self):
        self.section_links = []
        self.buffer_texts = []
        self.buffer_images = []
        self.section = False
        with self.app:
            self.channel_dialog = None
            self.chat_dialog = None
            for dialog in self.app.get_dialogs():
                if (dialog.chat.title or "").strip() == args.tgchannel.strip():
                    self.channel_dialog = dialog
                if (dialog.chat.title or "").strip() == args.tgchat.strip():
                    self.chat_dialog = dialog
                if self.channel_dialog is not None and self.chat_dialog is not None:
                    break
            if not self.channel_dialog:
                raise Exception("Channel not found, please check provided name")
            if not self.chat_dialog:
                raise Exception("Linked chat not found, please check provided name")
            self.channel_id = self.channel_dialog.chat.id
            self.chat_id = self.chat_dialog.chat.id
            for pair in self.structure:
                self.tg_process_element(pair)
            if self.buffer_texts or self.buffer_images:
                posts = self.split_to_messages(self.buffer_texts, self.buffer_images)
                self.post_wrapper(posts)
                self.buffer_texts = []
                self.buffer_images = []
            if not self.args.skip_until:
                navigation_text = [self.labels["general"]["general_impressions_text"]]
                if self.tg_heading:
                    navigation_text = [f"**{self.tg_heading}**", ""] + navigation_text
                for i, link in enumerate(self.section_links):
                    navigation_text.append(
                        f"{self.labels['general']['section']} {i + 1}: {link}"
                    )
                navigation_text = "\n".join(navigation_text)
                messages = self.post([(navigation_text.strip(), None)])
                if not self.args.dry_run:
                    self.app.pin_chat_message(
                        self.channel_id,
                        messages[0].id,
                        disable_notification=True,
                    )


class LatexExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.qcount = 0

    def tex_format_question(self, q):
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        res = (
            "\n\n\\begin{{minipage}}{{\\textwidth}}\\raggedright\n"
            "\\textbf{{Вопрос {}.}} {} \\newline".format(
                self.qcount if "number" not in q else q["number"],
                self.texyapper(q["question"]),
            )
        )
        if "number" not in q:
            self.qcount += 1
        res += "\n\\textbf{{Ответ: }}{} \\newline".format(self.texyapper(q["answer"]))
        if "zachet" in q:
            res += "\n\\textbf{{Зачёт: }}{} \\newline".format(
                self.texyapper(q["zachet"])
            )
        if "nezachet" in q:
            res += "\n\\textbf{{Незачёт: }}{} \\newline".format(
                self.texyapper(q["nezachet"])
            )
        if "comment" in q:
            res += "\n\\textbf{{Комментарий: }}{} \\newline".format(
                self.texyapper(q["comment"])
            )
        if "source" in q:
            res += "\n\\textbf{{Источник{}: }}{} \\newline".format(
                "и" if isinstance(q["source"], list) else "",
                self.texyapper(q["source"]),
            )
        if "author" in q:
            res += "\n\\textbf{{Автор: }}{} \\newline".format(
                self.texyapper(q["author"])
            )
        res += "\n\\end{minipage}\n"
        return res

    @staticmethod
    def texrepl(zz):
        zz = re.sub(r"{", r"\{", zz)
        zz = re.sub(r"}", r"\}", zz)
        zz = re.sub(r"\\(?![\}\{])", r"{\\textbackslash}", zz)
        zz = re.sub("%", "\\%", zz)
        zz = re.sub(r"\$", "\\$", zz)
        zz = re.sub("#", "\\#", zz)
        zz = re.sub("&", "\\&", zz)
        zz = re.sub("_", r"\_", zz)
        zz = re.sub(r"\^", r"{\\textasciicircum}", zz)
        zz = re.sub(r"\~", r"{\\textasciitilde}", zz)
        zz = re.sub(r'((\"(?=[ \.\,;\:\?!\)\]]))|("(?=\Z)))', "»", zz)
        zz = re.sub(r'(((?<=[ \.\,;\:\?!\(\[)])")|((?<=\A)"))', "«", zz)
        zz = re.sub('"', "''", zz)

        for match in sorted(
            [x for x in re_scaps.finditer(zz)],
            key=lambda x: len(x.group(2)),
            reverse=True,
        ):
            zz = zz.replace(match.group(2), "\\tsc{" + match.group(2).lower() + "}")

        torepl = [x.group(0) for x in re.finditer(re_url, zz)]
        for s in range(len(torepl)):
            item = torepl[s]
            while item[-1] in typotools.PUNCTUATION:
                item = item[:-1]
            while (
                item[-1] in typotools.CLOSING_BRACKETS
                and typotools.find_matching_opening_bracket(item, -1) is None
            ):
                item = item[:-1]
            while item[-1] in typotools.PUNCTUATION:
                item = item[:-1]
            torepl[s] = item
        torepl = sorted(set(torepl), key=len, reverse=True)
        hashurls = {}
        for s in torepl:
            hashurls[s] = hashlib.md5(s.encode("utf8")).hexdigest()
        for s in sorted(hashurls, key=len, reverse=True):
            zz = zz.replace(s, hashurls[s])
        hashurls = {v: k for k, v in hashurls.items()}
        for s in sorted(hashurls):
            zz = zz.replace(s, "\\url{{{}}}".format(hashurls[s].replace("\\\\", "\\")))

        zz = zz.replace(" — ", "{\\Hair}—{\\hair}")

        while "`" in zz:
            if zz.index("`") + 1 >= len(zz):
                zz = zz.replace("`", "")
            else:
                if zz.index("`") + 2 < len(zz) and re.search(
                    r"\s", zz[zz.index("`") + 2]
                ):
                    zz = zz[: zz.index("`") + 2] + "" + zz[zz.index("`") + 2 :]
                if zz.index("`") + 1 < len(zz) and re_lowercase.search(
                    zz[zz.index("`") + 1]
                ):
                    zz = (
                        zz[: zz.index("`") + 1]
                        + ""
                        + zz[zz.index("`") + 1]
                        + "\u0301"
                        + zz[zz.index("`") + 2 :]
                    )
                elif zz.index("`") + 1 < len(zz) and re_uppercase.search(
                    zz[zz.index("`") + 1]
                ):
                    zz = (
                        zz[: zz.index("`") + 1]
                        + ""
                        + zz[zz.index("`") + 1]
                        + "\u0301"
                        + zz[zz.index("`") + 2 :]
                    )
                zz = zz[: zz.index("`")] + zz[zz.index("`") + 1 :]

        return zz

    def texformat(self, s):
        res = ""
        for run in parse_4s_elem(s):
            if run[0] == "":
                res += self.texrepl(run[1])
            if run[0] == "screen":
                res += self.texrepl(run[1]["for_print"])
            if run[0] == "em":
                res += "\\emph{" + self.texrepl(run[1]) + "}"
            if run[0] == "img":
                parsed_image = parseimg(
                    run[1],
                    dimensions="ems",
                    tmp_dir=self.dir_kwargs.get("tmp_dir"),
                    targetdir=self.dir_kwargs.get("targetdir"),
                )
                imgfile = parsed_image["imgfile"]
                w = parsed_image["width"]
                h = parsed_image["height"]
                res += (
                    "\\includegraphics"
                    + "[width={}{}]".format(
                        "10em" if w == -1 else "{}em".format(w),
                        ", height={}em".format(h) if h != -1 else "",
                    )
                    + "{"
                    + imgfile
                    + "}"
                )
        while res.endswith("\n"):
            res = res[:-1]
        res = res.replace("\n", "  \\newline \n")
        return res

    def texyapper(self, e):
        if isinstance(e, str):
            return self.tex_element_layout(e)
        elif isinstance(e, list):
            if not any(isinstance(x, list) for x in e):
                return self.tex_element_layout(e)
            else:
                return "  \n".join([self.tex_element_layout(x) for x in e])

    def tex_element_layout(self, e):
        res = ""
        if isinstance(e, str):
            res = self.texformat(e)
            return res
        if isinstance(e, list):
            res = """
    \\begin{{compactenum}}
    {}
    \\end{{compactenum}}
    """.format(
                "\n".join(["\\item {}".format(self.tex_element_layout(x)) for x in e])
            )
        return res

    def export(self, outfilename):
        self.qcount = 1
        tex = """\\input{@header}\n\\begin{document}""".replace(
            "@header", os.path.basename(self.args.tex_header)
        )
        firsttour = True
        for element in self.structure:
            if element[0] == "heading":
                tex += "\n{{\\huge {}}}\n" "\\vspace{{0.8em}}\n".format(
                    self.tex_element_layout(element[1])
                )
            if element[0] == "date":
                tex += "\n{{\\large {}}}\n" "\\vspace{{0.8em}}\n".format(
                    self.tex_element_layout(element[1])
                )
            if element[0] in {"meta", "editor"}:
                tex += "\n{}\n\\vspace{{0.8em}}\n".format(
                    self.tex_element_layout(element[1])
                )
            elif element[0] == "section":
                tex += "\n{}\\section*{{{}}}\n\n".format(
                    "\\clearpage" if not firsttour else "",
                    self.tex_element_layout(element[1]),
                )
                firsttour = False
            elif element[0] == "Question":
                tex += self.tex_format_question(element[1])

        tex += "\\end{document}"

        with codecs.open(outfilename, "w", "utf8") as outfile:
            outfile.write(tex)
        cwd = os.getcwd()
        os.chdir(self.dir_kwargs["tmp_dir"])
        subprocess.call(
            shlex.split(
                'xelatex -synctex=1 -interaction=nonstopmode "{}"'.format(outfilename)
            )
        )
        targetdir = os.path.dirname(outfilename)
        os.chdir(cwd)
        pdf_filename = os.path.splitext(os.path.basename(outfilename))[0] + ".pdf"
        logger.info("Output: {}".format(os.path.join(targetdir, pdf_filename)))
        shutil.copy(os.path.join(self.dir_kwargs["tmp_dir"], pdf_filename), targetdir)
        if self.args.rawtex:
            shutil.copy(outfilename, targetdir)
            shutil.copy(self.args.tex_header, targetdir)
            shutil.copy(
                os.path.join(self.dir_kwargs["tmp_dir"], "fix-unnumbered-sections.sty"),
                targetdir,
            )


class DocxExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.qcount = 0

    def _docx_format(self, *args, **kwargs):
        kwargs.update(self.dir_kwargs)
        return self.docx_format(*args, **kwargs)

    def docx_format(self, el, para, whiten, **kwargs):
        if isinstance(el, list):

            if len(el) > 1 and isinstance(el[1], list):
                self.docx_format(el[0], para, whiten, **kwargs)
                licount = 0
                for li in el[1]:
                    licount += 1

                    para.add_run("\n{}. ".format(licount))
                    self.docx_format(li, para, whiten, **kwargs)
            else:
                licount = 0
                for li in el:
                    licount += 1

                    para.add_run("\n{}. ".format(licount))
                    self.docx_format(li, para, whiten, **kwargs)

        if isinstance(el, str):
            logger.debug("parsing element {}:".format(log_wrap(el)))

            if kwargs.get("remove_accents"):
                el = el.replace("\u0301", "")
            if kwargs.get("remove_brackets"):
                el = self.remove_square_brackets(el)
            else:
                el = replace_escaped(el)

            el = backtick_replace(el)

            for run in parse_4s_elem(el):
                if run[0] == "pagebreak":
                    if self.args.spoilers == "dots":
                        for _ in range(30):
                            para = self.doc.add_paragraph()
                            para.add_run(".")
                        para = self.doc.add_paragraph()
                    else:
                        para = self.doc.add_page_break()
                elif run[0] == "screen":
                    if kwargs.get("remove_accents") or kwargs.get("remove_brackets"):
                        r = para.add_run(replace_no_break_spaces(run[1]["for_screen"]))
                    else:
                        r = para.add_run(replace_no_break_spaces(run[1]["for_print"]))
                elif run[0] == "hyperlink" and not (
                    whiten and self.args.spoilers == "whiten"
                ):
                    r = self.add_hyperlink(para, run[1], run[1])
                elif run[0] == "img":
                    if run[1].endswith(".shtml"):
                        r = para.add_run(
                            "(ТУТ БЫЛА ССЫЛКА НА ПРОТУХШУЮ КАРТИНКУ)\n"
                        )  # TODO: добавить возможность пропускать кривые картинки опцией
                        continue
                    parsed_image = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=kwargs.get("tmp_dir"),
                        targetdir=kwargs.get("targetdir"),
                    )
                    imgfile = parsed_image["imgfile"]
                    width = parsed_image["width"]
                    height = parsed_image["height"]
                    r = para.add_run("\n")
                    try:
                        r.add_picture(
                            imgfile, width=Inches(width), height=Inches(height)
                        )
                    except UnrecognizedImageError:
                        sys.stderr.write(
                            f"python-docx can't recognize header for {imgfile}\n"
                        )
                    r.add_text("\n")
                    continue
                else:
                    r = para.add_run(replace_no_break_spaces(run[1]))
                    if run[0] == "em":
                        r.italic = True
                    elif run[0] == "sc":
                        r.small_caps = True
                    if whiten and self.args.spoilers == "whiten":
                        r.style = "Whitened"

    def add_hyperlink(self, paragraph, text, url):
        # adapted from https://github.com/python-openxml/python-docx/issues/610
        doc = self.doc
        run = paragraph.add_run(text)
        run.style = doc.styles["Hyperlink"]
        part = paragraph.part
        r_id = part.relate_to(
            url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True
        )
        hyperlink = docx.oxml.shared.OxmlElement("w:hyperlink")
        hyperlink.set(docx.oxml.shared.qn("r:id"), r_id)
        hyperlink.append(run._r)
        paragraph._p.append(hyperlink)
        return hyperlink

    def add_question(self, element, skip_qcount=False, screen_mode=False):
        q = element[1]
        p = self.doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(18)
        p.paragraph_format.keep_together = True
        if "number" not in q and not skip_qcount:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        p.add_run(
            "{question}. ".format(
                question=self.get_label(
                    q,
                    "question",
                    number=self.qcount if "number" not in q else q["number"],
                )
            )
        ).bold = True

        if "handout" in q:
            p.add_run("\n[{handout}: ".format(handout=self.get_label(q, "handout")))
            self._docx_format(
                q["handout"],
                p,
                WHITEN["handout"],
                remove_accents=screen_mode,
                remove_brackets=screen_mode,
            )
            p.add_run("\n]")
        if not self.args.noparagraph:
            p.add_run("\n")

        self._docx_format(
            q["question"],
            p,
            False,
            remove_accents=screen_mode,
            remove_brackets=screen_mode,
        )

        if not self.args.noanswers:
            if self.args.spoilers == "pagebreak":
                p = self.doc.add_page_break()
            elif self.args.spoilers == "dots":
                for _ in range(30):
                    p = self.doc.add_paragraph()
                    p.add_run(".")
                p = self.doc.add_paragraph()
            else:
                p = self.doc.add_paragraph()
            p.paragraph_format.keep_together = True
            p.paragraph_format.space_before = DocxPt(6)
            p.add_run(f"{self.get_label(q, 'answer')}: ").bold = True
            self._docx_format(q["answer"], p, True, remove_accents=screen_mode)

            for field in ["zachet", "nezachet", "comment", "source", "author"]:
                if field in q:
                    if field == "source":
                        p = self.doc.add_paragraph()
                        p.paragraph_format.keep_together = True
                    else:
                        p.add_run("\n")
                    p.add_run(f"{self.get_label(q, field)}: ").bold = True
                    self._docx_format(
                        q[field],
                        p,
                        WHITEN[field],
                        remove_accents=screen_mode,
                        remove_brackets=screen_mode,
                    )

    def export(self, outfilename):
        logger.debug(self.args.docx_template)
        self.doc = Document(self.args.docx_template)
        para = None
        logger.debug(log_wrap(self.structure))

        firsttour = True
        prev_element = None
        para = None
        page_break_before_heading = False
        for element in self.structure:
            if element[0] == "meta":
                para = self.doc.add_paragraph()
                if prev_element and prev_element[0] == "Question":
                    para.paragraph_format.space_before = DocxPt(18)
                self._docx_format(element[1], para, False)
                self.doc.add_paragraph()

            if element[0] in ["editor", "date", "heading", "section"]:
                if element[0] == "heading" and para is not None:
                    page_break_before_heading = True
                if para is None:
                    para = self.doc.paragraphs[0]
                    para.add_run(element[1])
                else:
                    para = self.doc.add_paragraph(element[1])
                if element[0] == "heading" and page_break_before_heading:
                    para.paragraph_format.page_break_before = True
                if element[0] == "section":
                    if not firsttour:
                        para.paragraph_format.page_break_before = True
                    else:
                        firsttour = False
                para.alignment = 1
                para.paragraph_format.keep_with_next = True
                para.add_run("\n")

            if element[0] == "Question":
                if self.args.screen_mode == "add_versions":
                    para = self.doc.add_paragraph()
                    para = self.doc.add_paragraph()
                    para.add_run("Версия для ведущего:").bold = True
                self.add_question(
                    element, screen_mode=self.args.screen_mode == "replace_all"
                )
                if self.args.screen_mode == "add_versions":
                    para = self.doc.add_paragraph()
                    para = self.doc.add_paragraph()
                    para.add_run("Версия для экрана:").bold = True
                    self.add_question(element, skip_qcount=True, screen_mode=True)
            prev_element = element

        self.doc.save(outfilename)
        logger.info("Output: {}".format(outfilename))


class PptxExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.config_path = os.path.abspath(self.args.pptx_config)
        with open(self.config_path, encoding="utf8") as f:
            self.c = toml.load(f)
        self.qcount = 0
        hs = self.labels["question_labels"]["handout"]
        self.re_handout_1 = re.compile(
            "\\[" + hs + ".(?P<body>.+?)\\]", flags=re.DOTALL
        )
        self.re_handout_2 = re.compile("^" + hs + ".(?P<body>.+?)$")

    def get_textbox_qnumber(self, slide):
        kwargs = {}
        for param in ("left", "top", "width", "height"):
            try:
                kwargs[param] = PptxInches(self.c["number_textbox"][param])
            except KeyError:
                pass

        return self.get_textbox(slide, **kwargs)

    def get_textbox(self, slide, left=None, top=None, width=None, height=None):
        if left is None:
            left = PptxInches(self.c["textbox"]["left"])
        if top is None:
            top = PptxInches(self.c["textbox"]["top"])
        if width is None:
            width = PptxInches(self.c["textbox"]["width"])
        if height is None:
            height = PptxInches(self.c["textbox"]["height"])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        return textbox

    def pptx_format(self, el, para, tf, slide):
        if isinstance(el, list):
            if len(el) > 1 and isinstance(el[1], list):
                self.pptx_format(el[0], para, tf, slide)
                licount = 0
                for li in el[1]:
                    licount += 1
                    r = para.add_run()
                    r.text = "\n{}. ".format(licount)
                    self.pptx_format(li, para, tf, slide)
            else:
                licount = 0
                for li in el:
                    licount += 1
                    r = para.add_run()
                    r.text = "\n{}. ".format(licount)
                    self.pptx_format(li, para, tf, slide)

        if isinstance(el, str):
            logger.debug("parsing element {}:".format(log_wrap(el)))
            el = backtick_replace(el)

            for run in parse_4s_elem(el):
                if run[0] in ("", "sc"):
                    r = para.add_run()
                    r.text = replace_no_break_spaces(run[1])

                elif run[0] == "screen":
                    r = para.add_run()
                    r.text = replace_no_break_spaces(run[1]["for_screen"])

                elif run[0] == "em":
                    r = para.add_run()
                    r.text = replace_no_break_spaces(run[1])
                    r.italic = True

                elif run[0] == "img":
                    pass  # image processing is moved to other places

    def pptx_process_text(self, s, image=None, strip_brackets=True):
        hs = self.labels["question_labels"]["handout_short"]
        if isinstance(s, list):
            for i in range(len(s)):
                s[i] = self.pptx_process_text(s[i], image=image)
            return s
        if not self.args.do_not_remove_accents:
            s = s.replace("\u0301", "")
        if strip_brackets:
            s = self.remove_square_brackets(s)
            s = s.replace("]\n", "]\n\n")
        else:
            s = replace_escaped(s)
        if image:
            s = re.sub("\\[" + hs + "(.+?)\\]", "", s, flags=re.DOTALL)
            s = s.strip()
        elif hs in s:
            re_hs = re.search("\\[" + hs + ".+?: ?(.+)\\]", s, flags=re.DOTALL)
            if re_hs:
                s = s.replace(re_hs.group(0), re_hs.group(1))
        s = re.sub(" +", " ", s)
        for punct in (".", ",", "!", "?", ":"):
            s = s.replace(" " + punct, punct)
        s = replace_no_break_spaces(s)
        s = s.strip()
        return s

    def _process_block(self, block):
        section = [x for x in block if x[0] == "section"]
        editor = [x for x in block if x[0] == "editor"]
        meta = [x for x in block if x[0] == "meta"]
        if not section and not editor and not meta:
            return
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True
        text_for_size = (
            (self.recursive_join([x[1] for x in section]) or "")
            + "\n"
            + (self.recursive_join([x[1] for x in editor]) or "")
            + "\n"
            + (self.recursive_join([x[1] for x in meta]) or "")
        )
        p = self.init_paragraph(tf, text=text_for_size)
        add_line_break = False
        if section:
            r = p.add_run()
            r.text = replace_no_break_spaces(self.pptx_process_text(section[0][1]))
            r.font.size = PptxPt(self.c["text_size_grid"]["section"])
            add_line_break = True
        if editor:
            r = p.add_run()
            r.text = replace_no_break_spaces(
                ("\n" if add_line_break else "")
                + self.pptx_process_text(editor[0][1])
                + "\n"
            )
            add_line_break = True
        if meta:
            for element in meta:
                r = p.add_run()
                r.text = replace_no_break_spaces(
                    ("\n" if add_line_break else "")
                    + self.pptx_process_text(element[1])
                    + "\n"
                )
                add_line_break = True

    def process_buffer(self, buffer):
        heading_block = []
        editor_block = []
        section_block = []
        block = heading_block
        for element in buffer:
            if element[0] == "section":
                block = section_block
            if element[0] == "editor" and not section_block:
                block = editor_block
            block.append(element)
        heading = [x for x in heading_block if x[0] == "heading"]
        ljheading = [x for x in heading_block if x[0] == "ljheading"]
        title_text = ljheading or heading
        date_text = [x for x in heading_block if x[0] == "date"]
        if title_text:
            if len(self.prs.slides) == 1:
                slide = self.prs.slides[0]
            else:
                slide = self.prs.slides.add_slide(self.TITLE_SLIDE)
            title = slide.shapes.title
            title.text = title_text[0][1]
            if date_text:
                subtitle = slide.placeholders[1]
                subtitle.text = date_text[0][1]
        for block in (editor_block, section_block):
            self._process_block(block)

    def set_question_number(self, slide, number):
        qntextbox = self.get_textbox_qnumber(slide)
        qtf = qntextbox.text_frame
        qtf_p = self.init_paragraph(qtf)
        qtf_r = qtf_p.add_run()
        qtf_r.text = number
        if self.c["number_textbox"].get("color"):
            qtf_r.font.color.rgb = RGBColor(*self.c["number_textbox"]["color"])

    def _get_handout_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                handout = self._get_handout_from_4s(el)
                if handout:
                    return handout
        elif isinstance(text, str):
            match_ = self.re_handout_1.search(text)
            if match_:
                return match_.group("body")
            else:
                lines = text.split("\n")
                for line in lines:
                    match_ = self.re_handout_2.search(line)
                    if match_:
                        return match_.group("body")

    def _get_image_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                image = self._get_image_from_4s(el)
                if image:
                    return image
        elif isinstance(text, str):
            for run in parse_4s_elem(text):
                if run[0] == "img":
                    parsed_image = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=self.dir_kwargs.get("tmp_dir"),
                        targetdir=self.dir_kwargs.get("targetdir"),
                    )
                    return parsed_image

    def make_slide_layout(self, image, slide, allowbigimage=True):
        if image:
            ratio = image["width"] / image["height"]
            img_base_width = PptxInches(image["width"])
            img_base_height = PptxInches(image["height"])
            base_left = PptxInches(self.c["textbox"]["left"])
            base_top = PptxInches(self.c["textbox"]["top"])
            base_width = PptxInches(self.c["textbox"]["width"])
            base_height = PptxInches(self.c["textbox"]["height"])
            big_mode = (
                image["big"] and not self.c.get("text_is_duplicated") and allowbigimage
            )
            if ratio < 1:  # vertical image
                max_width = base_width // 3
                if big_mode:
                    max_width *= 2
                if img_base_width > max_width or big_mode:
                    img_width = max_width
                    img_height = int(img_base_height * (max_width / img_base_width))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left + img_width
                top = base_top
                width = base_width - img_width
                height = base_height
                img_left = base_left
                img_top = int(base_top + 0.5 * (base_height - img_height))
            else:  # horizontal/square image
                max_height = base_height // 3
                if big_mode:
                    max_height *= 2
                if img_base_height > max_height or big_mode:
                    img_height = max_height
                    img_width = int(img_base_width * (max_height / img_base_height))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left
                top = base_top + img_height
                width = base_width
                height = base_height - img_height
                img_top = base_top
                img_left = int(base_left + 0.5 * (base_width - img_width))
            slide.shapes.add_picture(
                image["imgfile"],
                left=img_left,
                top=img_top,
                width=img_width,
                height=img_height,
            )
            textbox = slide.shapes.add_textbox(left, top, width, height)
            return textbox, (width * height) / (base_width * base_height)
        else:
            return self.get_textbox(slide), 1

    def add_slide_with_image(self, image, number=None):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        if number:
            self.set_question_number(slide, number)
        img_width = PptxInches(image["width"])
        img_height = PptxInches(image["height"])
        base_left = PptxInches(self.c["textbox"]["left"])
        base_top = PptxInches(self.c["textbox"]["top"])
        base_width = PptxInches(self.c["textbox"]["width"])
        base_height = PptxInches(self.c["textbox"]["height"])
        if image["big"] or img_width > base_width:
            img_width, img_height = (
                base_width,
                int(img_height * (base_width / img_width)),
            )
        if img_height > base_height:
            img_width, img_height = (
                int(img_width * (base_height / img_height)),
                base_height,
            )
        img_left = int(base_left + 0.5 * (base_width - img_width))
        img_top = int(base_top + 0.5 * (base_height - img_height))
        slide.shapes.add_picture(
            image["imgfile"],
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height,
        )

    def put_question_on_slide(self, image, slide, q, allowbigimage=True):
        textbox, coeff = self.make_slide_layout(
            image, slide, allowbigimage=allowbigimage
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        self.set_question_number(slide, self.number)
        question_text = self.pptx_process_text(q["question"], image=image)
        p = self.init_paragraph(tf, text=question_text, coeff=coeff)
        self.pptx_format(question_text, p, tf, slide)

    def recursive_join(self, s):
        if isinstance(s, str):
            return s
        if isinstance(s, list):
            return "\n".join(self.recursive_join(x) for x in s)

    def add_slide_with_handout(self, handout, number=None):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True
        if number is not None:
            self.set_question_number(slide, number)
        p = self.init_paragraph(tf, text=handout)
        r = p.add_run()
        r.text = self.pptx_process_text(handout)

    def process_question_text(self, q):
        image = self._get_image_from_4s(q["question"])
        handout = self._get_handout_from_4s(q["question"])
        if image:
            self.add_slide_with_image(image, number=self.number)
        elif handout:
            self.add_slide_with_handout(handout, number=self.number)
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        text_is_duplicated = bool(self.c.get("text_is_duplicated"))
        self.put_question_on_slide(
            image, slide, q, allowbigimage=not text_is_duplicated
        )
        if image and image["big"] and text_is_duplicated:
            self.add_slide_with_image(image, number=self.number)

    def process_question(self, q):
        if "number" not in q:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        self.number = str(self.qcount if "number" not in q else q["number"])

        if isinstance(q["question"], list):
            for i in range(len(q["question"][1])):
                qn = copy.deepcopy(q)
                qn["question"][1] = q["question"][1][: i + 1]
                self.process_question_text(qn)
        else:
            self.process_question_text(q)

        if self.c["add_plug"]:
            slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
            self.set_question_number(slide, self.number)
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        self.set_question_number(slide, self.number)
        fields = ["answer"]
        if q.get("zachet") and self.c.get("add_zachet"):
            fields.append("zachet")
        if self.c["add_comment"] and "comment" in q:
            fields.append("comment")
        textbox = None
        coeff = 1
        for field in fields:
            image = self._get_image_from_4s(q[field])
            if image:
                textbox, coeff = self.make_slide_layout(image, slide)
                break
        if not textbox:
            textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True

        text_for_size = self.recursive_join(
            self.pptx_process_text(q["answer"], strip_brackets=False)
        )
        if q.get("zachet") and self.c.get("add_zachet"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["zachet"], strip_brackets=False)
            )
        if q.get("comment") and self.c.get("add_comment"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["comment"])
            )
        p = self.init_paragraph(tf, text=text_for_size, coeff=coeff)
        r = p.add_run()
        r.text = f"{self.get_label(q, 'answer')}: "
        r.font.bold = True
        self.pptx_format(
            self.pptx_process_text(q["answer"], strip_brackets=False), p, tf, slide
        )
        if q.get("zachet") and self.c.get("add_zachet"):
            zachet_text = self.pptx_process_text(q["zachet"], strip_brackets=False)
            r = p.add_run()
            r.text = f"\n{self.get_label(q, 'zachet')}: "
            r.font.bold = True
            self.pptx_format(zachet_text, p, tf, slide)
        if self.c["add_comment"] and "comment" in q:
            comment_text = self.pptx_process_text(q["comment"])
            r = p.add_run()
            r.text = f"\n{self.get_label(q, 'comment')}: "
            r.font.bold = True
            self.pptx_format(comment_text, p, tf, slide)

    def determine_size(self, text, coeff=1):
        text = self.recursive_join(text)
        len_for_size = round((len(text) + 50 * text.count("\n")) / coeff)
        for element in self.c["text_size_grid"]["elements"]:
            if len_for_size <= element["length"]:
                return element["size"]
        return self.c["text_size_grid"]["smallest"]

    def init_paragraph(self, text_frame, text=None, coeff=1):
        p = text_frame.paragraphs[0]
        p.font.name = self.c["font"]["name"]
        size = self.c["text_size_grid"]["default"]
        if text:
            size = self.determine_size(text, coeff=coeff)
        p.font.size = PptxPt(size)
        return p

    def export(self, outfilename):
        self.outfilename = outfilename
        wd = os.getcwd()
        os.chdir(os.path.dirname(self.config_path))
        template = os.path.abspath(self.c["template_path"])
        os.chdir(wd)
        self.prs = Presentation(template)
        self.TITLE_SLIDE = self.prs.slide_layouts[0]
        self.BLANK_SLIDE = self.prs.slide_layouts[6]
        buffer = []
        for element in self.structure:
            if element[0] != "Question":
                buffer.append(element)
                continue
            if element[0] == "Question":
                if buffer:
                    self.process_buffer(buffer)
                    buffer = []
                self.process_question(element[1])
        self.prs.save(outfilename)
        logger.info("Output: {}".format(outfilename))


class StatsAdder(BaseExporter):
    @staticmethod
    def patch_question(question, message):
        if "comment" not in question:
            question["comment"] = message
        elif isinstance(question["comment"], str):
            question["comment"] += "\n" + message
        elif isinstance(question["comment"], list):
            if len(question["comment"]) > 1:
                if isinstance(question["comment"][1], list):
                    question["comment"][1].append(message)
                else:
                    question["comment"].append(message)
            else:
                question["comment"].append(message)

    @staticmethod
    def get_tournament_results(id_):
        req = requests.get(
            "https://api.rating.chgk.net"
            + f"/tournaments/{id_}/results.json"
            + "?includeMasksAndControversials=1"
        )
        return req.json()

    def process_tournament(self, results):
        for res in results:
            if not res.get("mask"):
                continue
            self.total_teams += 1
            name = res["current"]["name"]
            mask = list(res["mask"])
            if self.args.question_range:
                start, end = self.args.question_range.split("-")
                start = int(start)
                end = int(end)
            else:
                start = 0
                end = 9999
            qnum = 1
            for i, q in enumerate(mask):
                if not start <= (i + 1) <= end:
                    continue
                if q == "1":
                    self.q_counter[qnum] += 1
                    self.q_to_teams[qnum].add(name)
                qnum += 1

    def export(self, outfilename):
        self.q_to_teams = defaultdict(set)
        self.total_teams = 0
        self.q_counter = Counter()
        if self.args.rating_ids:
            ids = [x.strip() for x in self.args.rating_ids.split(",") if x.strip()]
            for id_ in ids:
                results = self.get_tournament_results(id_)
                self.process_tournament(results)
        elif self.args.custom_csv:
            results = custom_csv_to_results(self.args.custom_csv)
            self.process_tournament(results)
        qnumber = 1
        for element in self.structure:
            if element[0] != "Question" or str(element[1].get("number")).startswith(
                "0"
            ):
                continue
            scored_teams = self.q_counter[qnumber]
            label = self.labels["general"]["right_answers_for_stats"]
            share = scored_teams / self.total_teams
            message = (
                f"{label}: {scored_teams}/{self.total_teams} ({round(share * 100)}%)"
            )
            if scored_teams > 0 and scored_teams <= self.args.team_naming_threshold:
                teams = ", ".join(sorted(self.q_to_teams[qnumber]))
                message += f" ({teams})"
            self.patch_question(element[1], message)
            qnumber += 1
        with open(outfilename, "w", encoding="utf8") as f:
            f.write(compose_4s(self.structure, args=self.args))
            logger.info(f"Output: {outfilename}")


class Imgur:
    def __init__(self, client_id):
        self.client_id = client_id
        self.cache_file_path = os.path.join(get_chgksuite_dir(), "image_cache.json")
        if os.path.isfile(self.cache_file_path):
            try:
                with open(self.cache_file_path) as f:
                    self.cache = json.load(f)
            except json.decoder.JSONDecodeError:
                self.cache = {}
        else:
            self.cache = {}

    def upload_image(self, path, title=None):
        with open(path, "rb") as image_file:
            binary_data = image_file.read()
        image_bytes = base64.b64encode(binary_data)
        image = image_bytes.decode("utf8", errors="replace")
        sha256 = hashlib.sha256(image_bytes).hexdigest()
        if sha256 in self.cache:
            return {"data": {"link": self.cache[sha256]}}
        payload = {
            "album_id": None,
            "image": image,
            "title": title,
            "description": None,
        }
        retries = 0
        req = None
        while (not req or req.status_code != 200) and retries < 10:
            req = requests.post(
                "https://api.imgur.com/3/image",
                json=payload,
                headers={"Authorization": f"Client-ID {self.client_id}"},
            )
            if req.status_code != 200:
                sys.stderr.write(f"got 403 from imgur, retry {retries + 1}...")
                retries += 1
                time.sleep(5)
        try:
            assert req.status_code == 200
            json_ = req.json()
            self.cache[sha256] = json_["data"]["link"]
            with open(self.cache_file_path, "w", encoding="utf8") as f:
                json.dump(self.cache, f, indent=2, sort_keys=True)
            return json_
        except Exception as e:
            raise Exception(
                f"Imgur API error code {req.status_code}: "
                f"{req.content.decode('utf8', errors='replace')}, raw exception data: "
                f"{type(e)} {e}"
            )


class LjExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.lj = ServerProxy("http://www.livejournal.com/interface/xmlrpc").LJ.XMLRPC
        self.im = Imgur(self.args.imgur_client_id or IMGUR_CLIENT_ID)

    def get_chal(self):
        chal = None
        chal = retry_wrapper(self.lj.getchallenge)["challenge"]
        response = md5(
            chal.encode("utf8") + md5(self.args.password.encode("utf8")).encode("utf8")
        )
        return (chal, response)

    def split_into_tours(self):
        general_impression = self.args.genimp
        result = []
        current = []
        mode = "meta"
        for _, element in enumerate(self.structure):
            if element[0] != "Question":
                if mode == "meta":
                    current.append(element)
                elif element[0] == "section":
                    result.append(current)
                    current = [element]
                    mode = "meta"
                else:
                    current.append(element)
            else:
                if mode == "meta":
                    current.append(element)
                    mode = "questions"
                else:
                    current.append(element)
        result.append(current)
        globalheading = find_heading(result[0])[1][1]
        globalsep = "." if not globalheading.endswith(".") else ""
        currentheading = result[0][find_heading(result[0])[0]][1]
        result[0][find_heading(result[0])[0]][1] += "{} {}".format(
            "." if not currentheading.endswith(".") else "", find_tour(result[0])[1][1]
        )
        for tour in result[1:]:
            if not find_heading(tour):
                tour.insert(
                    0,
                    [
                        "ljheading",
                        "{}{} {}".format(
                            globalheading, globalsep, find_tour(tour)[1][1]
                        ),
                    ],
                )
        if general_impression:
            result.append(
                [
                    [
                        "ljheading",
                        "{}{} {}".format(
                            globalheading,
                            globalsep,
                            self.labels["general"]["general_impressions_caption"],
                        ),
                    ],
                    ["meta", self.labels["general"]["general_impressions_text"]],
                ]
            )
        return result

    def _lj_post(self, stru, edit=False, add_params=None):

        now = datetime.datetime.now()
        year = now.strftime("%Y")
        month = now.strftime("%m")
        day = now.strftime("%d")
        hour = now.strftime("%H")
        minute = now.strftime("%M")

        chal, response = self.get_chal()

        params = {
            "username": self.args.login,
            "auth_method": "challenge",
            "auth_challenge": chal,
            "auth_response": response,
            "subject": stru["header"],
            "event": stru["content"],
            "year": year,
            "mon": month,
            "day": day,
            "hour": hour,
            "min": minute,
        }
        if edit:
            params["itemid"] = stru["itemid"]
        if add_params:
            params.update(add_params)

        try:
            post = retry_wrapper(
                self.lj.editevent if edit else self.lj.postevent, [params]
            )
            logger.info("Edited a post" if edit else "Created a post")
            logger.debug(log_wrap(post))
            time.sleep(5)
        except Exception as e:
            sys.stderr.write(
                "Error issued by LJ API: {}".format(traceback.format_exc(e))
            )
            sys.exit(1)
        return post

    def _lj_comment(self, stru):
        chal, response = self.get_chal()
        params = {
            "username": self.args.login,
            "auth_method": "challenge",
            "auth_challenge": chal,
            "auth_response": response,
            "journal": stru["journal"],
            "ditemid": stru["ditemid"],
            "parenttalkid": 0,
            "body": stru["content"],
            "subject": stru["header"],
        }
        try:
            comment = retry_wrapper(self.lj.addcomment, [params])
        except Exception as e:
            sys.stderr.write(
                "Error issued by LJ API: {}".format(traceback.format_exc(e))
            )
            sys.exit(1)
        logger.info("Added a comment")
        logger.debug(log_wrap(comment))
        time.sleep(random.randint(7, 12))

    def lj_post(self, stru, edit=False):

        add_params = {}
        community = self.args.community
        if community:
            add_params["usejournal"] = community
        elif self.args.security == "public":
            pass
        elif self.args.security:
            add_params["security"] = "usemask"
            add_params["allowmask"] = (
                "1" if self.args.security == "friends" else self.args.security
            )
        else:
            add_params["security"] = "private"

        journal = community if community else self.args.login

        post = self._lj_post(stru[0], edit=edit, add_params=add_params)

        comments = stru[1:]

        if not comments:
            return post

        for comment in stru[1:]:
            comment["ditemid"] = post["ditemid"]
            comment["journal"] = journal
            self._lj_comment(comment)

        return post

    def lj_process(self, structure):
        final_structure = [{"header": "", "content": ""}]
        i = 0
        heading = ""
        ljheading = ""

        def yapper(x):
            return self.htmlyapper(x)

        while i < len(structure) and structure[i][0] != "Question":
            if structure[i][0] == "heading":
                final_structure[0]["content"] += "<center>{}</center>".format(
                    yapper(structure[i][1])
                )
                heading = yapper(structure[i][1])
            if structure[i][0] == "ljheading":
                # final_structure[0]['header'] = structure[i][1]
                ljheading = yapper(structure[i][1])
            if structure[i][0] == "date":
                final_structure[0]["content"] += "\n<center>{}</center>".format(
                    yapper(structure[i][1])
                )
            if structure[i][0] == "editor":
                final_structure[0]["content"] += "\n<center>{}</center>".format(
                    yapper(structure[i][1])
                )
            if structure[i][0] == "meta":
                final_structure[0]["content"] += "\n{}".format(yapper(structure[i][1]))
            i += 1

        if ljheading != "":
            final_structure[0]["header"] = ljheading
        else:
            final_structure[0]["header"] = heading

        for element in structure[i:]:
            if element[0] == "Question":
                formatted = self.get_label(element[1], "question", number=self.counter)
                final_structure.append(
                    {
                        "header": formatted,
                        "content": self.html_format_question(element[1]),
                    }
                )
                self.counter += 1
            if element[0] == "meta":
                final_structure.append({"header": "", "content": yapper(element[1])})

        if not final_structure[0]["content"]:
            final_structure[0]["content"] = self.labels["general"][
                "general_impressions_text"
            ]
        if self.args.debug:
            with codecs.open("lj.debug", "w", "utf8") as f:
                f.write(log_wrap(final_structure))
        return final_structure

    def htmlyapper(self, e):
        if isinstance(e, str):
            return self.html_element_layout(e)
        elif isinstance(e, list):
            if not any(isinstance(x, list) for x in e):
                return self.html_element_layout(e)
            else:
                return "\n".join([self.html_element_layout(x) for x in e])

    def html_element_layout(self, e):
        res = ""
        if isinstance(e, str):
            res = self.htmlformat(e)
            return res
        if isinstance(e, list):
            res = "\n".join(
                [
                    "{}. {}".format(en + 1, self.html_element_layout(x))
                    for en, x in enumerate(e)
                ]
            )
            return res

    def html_format_question(self, q):
        def yapper(x):
            return self.htmlyapper(x)

        if "setcounter" in q:
            self.counter = int(q["setcounter"])
        res = "<strong>{question}.</strong> {content}".format(
            question=self.get_label(q, "question", self.counter),
            content=yapper(q["question"])
            + ("\n<lj-spoiler>" if not args.nospoilers else ""),
        )
        if "number" not in q:
            self.counter += 1
        for field in ("answer", "zachet", "nezachet", "comment", "source", "author"):
            if field in q:
                res += "\n<strong>{field}: </strong>{content}".format(
                    field=self.get_label(q, field), content=yapper(q[field])
                )
        if not args.nospoilers:
            res += "</lj-spoiler>"
        return res

    @staticmethod
    def htmlrepl(zz):
        zz = zz.replace("&", "&amp;")
        zz = zz.replace("<", "&lt;")
        zz = zz.replace(">", "&gt;")

        while "`" in zz:
            if zz.index("`") + 1 >= len(zz):
                zz = zz.replace("`", "")
            else:
                if zz.index("`") + 2 < len(zz) and re.search(
                    r"\s", zz[zz.index("`") + 2]
                ):
                    zz = zz[: zz.index("`") + 2] + "" + zz[zz.index("`") + 2 :]
                if zz.index("`") + 1 < len(zz) and re_lowercase.search(
                    zz[zz.index("`") + 1]
                ):
                    zz = (
                        zz[: zz.index("`") + 1]
                        + ""
                        + zz[zz.index("`") + 1]
                        + "&#x0301;"
                        + zz[zz.index("`") + 2 :]
                    )
                elif zz.index("`") + 1 < len(zz) and re_uppercase.search(
                    zz[zz.index("`") + 1]
                ):
                    zz = (
                        zz[: zz.index("`") + 1]
                        + ""
                        + zz[zz.index("`") + 1]
                        + "&#x0301;"
                        + zz[zz.index("`") + 2 :]
                    )
                zz = zz[: zz.index("`")] + zz[zz.index("`") + 1 :]

        return zz

    def htmlformat(self, s):
        res = ""
        for run in parse_4s_elem(s):
            if run[0] in ("", "hyperlink"):
                res += self.htmlrepl(run[1])
            if run[0] == "screen":
                res += self.htmlrepl(run[1]["for_screen"])
            if run[0] == "em":
                res += "<em>" + self.htmlrepl(run[1]) + "</em>"
            if run[0] == "img":
                parsed_image = parseimg(
                    run[1],
                    dimensions="pixels",
                    targetdir=self.dir_kwargs.get("targetdir"),
                    tmp_dir=self.dir_kwargs.get("tmp_dir"),
                )
                imgfile = parsed_image["imgfile"]
                w = parsed_image["width"]
                h = parsed_image["height"]
                if os.path.isfile(imgfile):
                    uploaded_image = self.im.upload_image(imgfile, title=imgfile)
                    imgfile = uploaded_image["data"]["link"]

                res += '<img{}{} src="{}"/>'.format(
                    "" if w == -1 else " width={}".format(w),
                    "" if h == -1 else " height={}".format(h),
                    imgfile,
                )
        res = replace_no_break_spaces(res)
        return res

    def export(self):
        args = self.args
        if not args.community:
            args.community = ""
        if not args.login:
            print("Login not specified.")
            sys.exit(1)
        elif not args.password:
            import getpass

            args.password = getpass.getpass()

        self.counter = 1
        if args.splittours:
            tours = self.split_into_tours()
            strus = []
            for tour in tours:
                stru = self.lj_process(tour)
                post = self.lj_post(stru)
                strus.append((stru, post))
            if args.navigation:
                navigation = generate_navigation(strus)
                for i, (stru, post) in enumerate(strus):
                    newstru = {
                        "header": stru[0]["header"],
                        "content": stru[0]["content"] + "\n\n" + navigation[i],
                        "itemid": post["itemid"],
                    }
                    self.lj_post([newstru], edit=True)
        else:
            stru = self.lj_process(self.structure)
            post = self.lj_post(stru)


def process_file(filename, tmp_dir, sourcedir, targetdir):
    global args
    dir_kwargs = dict(tmp_dir=tmp_dir, targetdir=targetdir)

    if isinstance(filename, list):
        structure = []
        for x in filename:
            structure.extend(parse_filepath(os.path.join(targetdir, x)))
        filename = make_merged_filename(filename)
    else:
        structure = parse_filepath(os.path.join(targetdir, filename))

    if args.debug:
        debug_fn = os.path.join(
            targetdir,
            make_filename(os.path.basename(filename), "dbg", args),
        )
        with codecs.open(debug_fn, "w", "utf8") as output_file:
            output_file.write(json.dumps(structure, indent=2, ensure_ascii=False))

    if not args.filetype:
        print("Filetype not specified.")
        sys.exit(1)
    if args.filetype == "docx":
        spoilers = args.spoilers
    else:
        spoilers = "off" if args.nospoilers else "on"
    logger.info("Exporting to {}, spoilers are {}...\n".format(args.filetype, spoilers))

    if args.filetype == "docx":

        if args.screen_mode == "off":
            addsuffix = ""
        elif args.screen_mode == "replace_all":
            addsuffix = "_screen"
        elif args.screen_mode == "add_versions":
            addsuffix = "_screen_versions"
        if args.spoilers != "off":
            addsuffix += "_spoilers"
        outfilename = os.path.join(
            targetdir, make_filename(filename, "docx", args, addsuffix=addsuffix)
        )
        exporter = DocxExporter(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if args.filetype == "tex":
        outfilename = os.path.join(tmp_dir, make_filename(filename, "tex", args))
        exporter = LatexExporter(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if args.filetype == "lj":
        exporter = LjExporter(structure, args, dir_kwargs)
        exporter.export()

    if args.filetype == "base":
        exporter = DbExporter(structure, args, dir_kwargs)
        outfilename = os.path.join(targetdir, make_filename(filename, "txt", args))
        exporter.export(outfilename)

    if args.filetype == "redditmd":
        exporter = RedditExporter(structure, args, dir_kwargs)
        outfilename = os.path.join(targetdir, make_filename(filename, "md", args))
        exporter.export(outfilename)

    if args.filetype == "pptx":
        outfilename = os.path.join(targetdir, make_filename(filename, "pptx", args))
        exporter = PptxExporter(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if args.filetype == "add_stats":
        outfilename = os.path.join(
            targetdir,
            make_filename(filename, "4s", args, addsuffix="_with_stats"),
        )
        exporter = StatsAdder(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if args.filetype == "telegram":
        exporter = TelegramExporter(structure, args, dir_kwargs)
        exporter.export()

    if not console_mode:
        input("Press Enter to continue...")


def main():
    print("This program was not designed to run standalone.")
    input("Press Enter to continue...")


if __name__ == "__main__":
    main()
